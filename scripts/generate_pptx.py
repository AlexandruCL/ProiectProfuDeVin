#!/usr/bin/env python
# -*- coding: utf-8 -*-

from __future__ import annotations

from pathlib import Path
from typing import Sequence

try:
    from pptx import Presentation
except ImportError as exc:
    raise SystemExit(
        "Missing dependency: python-pptx. Install with: pip install python-pptx"
    ) from exc


def add_title_slide(prs: Presentation, title: str, subtitle: str) -> None:
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = title
    slide.placeholders[1].text = subtitle


def add_bullets_slide(prs: Presentation, title: str, bullets: Sequence[str]) -> None:
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = title
    tf = slide.placeholders[1].text_frame
    tf.clear()
    for i, b in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = b
        p.level = 0


def main() -> None:
    prs = Presentation()

    add_title_slide(
        prs,
        "ProiectProfuDeVin",
        "Django wines & spirits — Render + Supabase deployment",
    )

    add_bullets_slide(
        prs,
        "Goal",
        [
            "Deploy Django reliably with Docker",
            "Use Supabase Postgres for persistence",
            "Handle static + media assets",
            "Seed product data automatically",
        ],
    )

    add_bullets_slide(
        prs,
        "Architecture",
        [
            "Render Web Service: Docker + Gunicorn",
            "Supabase Postgres: DATABASE_URL",
            "Supabase Storage (recommended) for media",
        ],
    )

    add_bullets_slide(
        prs,
        "Startup Flow",
        [
            "migrate",
            "optional: RUN_SEED (import wines/spirits)",
            "optional: CREATE_SUPERUSER",
            "optional: UPLOAD_MEDIA",
            "start Gunicorn",
        ],
    )

    add_bullets_slide(
        prs,
        "Key Env Vars",
        [
            "DATABASE_URL, SECRET_KEY, DEBUG=False",
            "RUN_SEED / CREATE_SUPERUSER / UPLOAD_MEDIA",
            "SUPABASE_STORAGE_ENABLED + project ref + bucket + S3 keys",
        ],
    )

    add_bullets_slide(
        prs,
        "Troubleshooting",
        [
            "Tables but no data: DATABASE_URL missing/wrong",
            "Import skipped: RUN_SEED not enabled",
            "Images broken: bucket not public or URL mismatch",
        ],
    )

    out_dir = Path(__file__).resolve().parents[1] / "presentation"
    out_dir.mkdir(parents=True, exist_ok=True)
    out_path = out_dir / "ProiectProfuDeVin_Deployment.pptx"
    prs.save(out_path)
    print(f"Wrote {out_path}")


if __name__ == "__main__":
    main()
